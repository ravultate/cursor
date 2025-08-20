import streamlit as st
import os
import logging
from typing import List, Dict, Any, Optional
from datetime import datetime
import json
import time
from openai import AzureOpenAI
from azure.search.documents import SearchClient
from azure.search.documents.models import VectorizedQuery
from azure.core.credentials import AzureKeyCredential
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import base64
import pandas as pd
from docx import Document
from docx.shared import Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from azure.identity import DefaultAzureCredential
from azure.core.credentials import AzureKeyCredential
from auth_utils import get_search_credential,get_openai_client,get_content_safety_client
from safety_utils import is_text_safe, enforce_output_safety, groundedness_check

load_dotenv()

# Environment variables
OPENAI_ENDPOINT_URL = os.getenv("OPENAI_ENDPOINT_URL")
DEPLOYMENT_NAME = os.getenv("DEPLOYMENT_NAME")
EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL", "text-embedding-3-large")
SEARCH_SERVICE_ENDPOINT = os.getenv("SEARCH_SERVICE_ENDPOINT")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

logger = logging.getLogger(__name__)

class FinancialInsightsAgent:
    """Agent specialized in generating financial due diligence insights"""
    
    def __init__(self, search_client: SearchClient, openai_client: AzureOpenAI, deployment_name: str):
        self.search_client = search_client
        self.openai_client = openai_client
        self.deployment_name = deployment_name
    
    def generate_query_embedding(self, query: str) -> List[float]:
        """Generate embedding for the query"""
        try:
            response = self.openai_client.embeddings.create(
                input=query,
                model=EMBEDDING_MODEL
            )
            return response.data[0].embedding
        except Exception as e:
            logger.error(f"Error generating query embedding: {str(e)}")
            return [0.0] * 3072
    
    def search_financial_data(self, query: str, k: int = 25) -> List[Dict]:
        """Search for financial data with enhanced parameters"""
        try:
            query_vector = self.generate_query_embedding(query)
            
            vector_query = VectorizedQuery(
                vector=query_vector,
                k_nearest_neighbors=k,
                fields="content_vector"
            )
            
            search_params = {
                "search_text": query,
                "vector_queries": [vector_query],
                "select": ["id", "content", "page_number", "file_name", "chunk_id"],
                "top": k,
                "query_type": "semantic",
                "semantic_configuration_name": "semantic_config",
                "query_caption": "extractive"
            }
            
            results = self.search_client.search(**search_params)
            
            search_results = []
            for result in results:
                doc_result = {
                    "id": result.get("id"),
                    "content": result.get("content"),
                    "page_number": result.get("page_number"),
                    "file_name": result.get("file_name"),
                    "chunk_id": result.get("chunk_id"),
                    "search_score": result.get("@search.score", 0)
                }
                search_results.append(doc_result)
            
            return search_results
            
        except Exception as e:
            logger.error(f"Error in financial data search: {str(e)}")
            return []
    
    def get_analysis_specific_queries(self, analysis_type: str, company_name: str) -> List[str]:
        """Get analysis-specific search queries that are truly unique for each type"""
        
        query_templates = {
            "comprehensive": [
                f"{company_name} financial statements income statement balance sheet",
                f"{company_name} revenue growth profit margins EBITDA operating income",
                f"{company_name} assets liabilities equity debt structure",
                f"{company_name} cash flow operating investing financing activities",
                f"{company_name} market share competitive position industry analysis",
                f"{company_name} risk factors operational financial market risks",
                f"{company_name} management team leadership governance structure",
                f"{company_name} business model strategy competitive advantages",
                f"{company_name} regulatory compliance legal issues",
                f"{company_name} ESG factors sustainability environmental social governance"
            ],
            "financial_performance": [
                f"{company_name} revenue growth trends year over year quarterly",
                f"{company_name} profitability analysis gross operating net margins",
                f"{company_name} EBITDA EBIT operating income trends",
                f"{company_name} return on equity ROE return on assets ROA",
                f"{company_name} working capital management efficiency ratios",
                f"{company_name} cost structure analysis operating expenses",
                f"{company_name} financial ratios liquidity profitability efficiency",
                f"{company_name} earnings per share EPS growth dividend policy"
            ],
            "risk_assessment": [
                f"{company_name} credit risk financial stability debt coverage",
                f"{company_name} market risk exposure currency interest rate",
                f"{company_name} operational risk business continuity",
                f"{company_name} regulatory compliance risk legal exposure",
                f"{company_name} cybersecurity risk data protection",
                f"{company_name} supply chain risk dependencies vulnerabilities",
                f"{company_name} ESG risk environmental social governance",
                f"{company_name} concentration risk customer supplier geographic"
            ],
            "market_analysis": [
                f"{company_name} industry analysis market trends competitive landscape",
                f"{company_name} market share competitive position ranking",
                f"{company_name} customer analysis demographics retention satisfaction",
                f"{company_name} pricing strategy competitive pricing analysis",
                f"{company_name} product portfolio market positioning",
                f"{company_name} geographic presence market expansion",
                f"{company_name} distribution channels sales network",
                f"{company_name} innovation R&D technology competitive advantage"
            ],
            "valuation": [
                f"{company_name} valuation multiples price earnings P/E ratio",
                f"{company_name} enterprise value EV EBITDA multiples",
                f"{company_name} book value market value price to book",
                f"{company_name} comparable company analysis peer valuation",
                f"{company_name} discounted cash flow DCF valuation model",
                f"{company_name} asset valuation tangible intangible assets",
                f"{company_name} market capitalization equity value",
                f"{company_name} transaction multiples M&A comparables"
            ]
        }
        
        return query_templates.get(analysis_type, query_templates["comprehensive"])
    
    def generate_comprehensive_insights(self, company_name: str, analysis_type: str, 
                                      custom_focus: str = "") -> Dict[str, Any]:
        """Generate comprehensive financial insights for due diligence"""
        
        # Get analysis-specific queries
        queries = self.get_analysis_specific_queries(analysis_type, company_name)
        if custom_focus:
            queries.append(f"{company_name} {custom_focus}")
        
        # Search for relevant documents
        all_results = []
        for query in queries:
            results = self.search_financial_data(query, k=8)
            all_results.extend(results)
        
        # Remove duplicates and get top results
        unique_results = {}
        for result in all_results:
            doc_id = result.get("id")
            if doc_id not in unique_results:
                unique_results[doc_id] = result
        
        top_results = sorted(unique_results.values(), 
                           key=lambda x: x.get("search_score", 0), 
                           reverse=True)[:20]
        
        # Generate comprehensive insights
        insights = self.generate_structured_insights(company_name, analysis_type, top_results)
        
        return {
            "company_name": company_name,
            "analysis_type": analysis_type,
            "insights": insights,
            "source_documents": top_results,
            "generation_timestamp": datetime.now().isoformat()
        }
    
    def generate_structured_insights(self, company_name: str, analysis_type: str, 
                                   documents: List[Dict]) -> Dict[str, Any]:
        """Generate structured insights based on analysis type"""
        
        # Prepare context
        context_parts = []
        for doc in documents:
            content = doc.get("content", "")
            file_name = doc.get("file_name", "Unknown")
            page_num = doc.get("page_number", "Unknown")
            context_parts.append(f"Source: {file_name} (Page {page_num})\n{content}")
        
        context_text = "\n\n".join(context_parts)
        
        # Generate analysis-specific prompts
        prompts = self.get_analysis_specific_prompts(analysis_type)
        
        insights_prompt = f"""
        You are a senior financial analyst at a Big 4 consulting firm preparing a comprehensive due diligence report.
        
        Company: {company_name}
        Analysis Focus: {analysis_type}
        
        {prompts['instructions']}
        
        Based on the provided financial documents, generate a comprehensive analysis following this structure:
        {prompts['structure']}
        
        Requirements:
        - Provide specific numerical data and metrics when available
        - Include professional risk-adjusted perspectives
        - Make actionable recommendations for senior management
        - Use Big 4 consulting terminology and frameworks
        - Ensure each section is comprehensive and detailed
        
        Context from Documents:
        {context_text}
        
        Provide a JSON response with the following structure:
        {prompts['json_structure']}
        """
        
        try:
            response = self.openai_client.chat.completions.create(
                model=self.deployment_name,
                messages=[{"role": "user", "content": insights_prompt}],
                max_tokens=4000,
                temperature=0.1
            )
            
            content = response.choices[0].message.content
            
            # Extract JSON from response
            import re
            json_match = re.search(r'\{.*\}', content, re.DOTALL)
            if json_match:
                return json.loads(json_match.group())
            else:
                return {"error": "Could not parse insights response"}
                
        except Exception as e:
            logger.error(f"Error generating structured insights: {str(e)}")
            return {"error": f"Error generating insights: {str(e)}"}
    
    def get_analysis_specific_prompts(self, analysis_type: str) -> Dict[str, str]:
        """Get analysis-specific prompts for different types"""
        
        prompts = {
            "comprehensive": {
                "instructions": """
                Conduct a comprehensive financial due diligence analysis covering all aspects of the company's financial health, 
                market position, operational efficiency, and strategic outlook. This analysis should be suitable for M&A transactions, 
                investment decisions, or strategic planning.
                """,
                "structure": """
                1. EXECUTIVE SUMMARY (Key findings and investment thesis)
                2. FINANCIAL PERFORMANCE ANALYSIS (Revenue, profitability, margins, trends)
                3. FINANCIAL POSITION ANALYSIS (Balance sheet strength, debt, liquidity)
                4. CASH FLOW ANALYSIS (Operating, investing, financing cash flows)
                5. MARKET POSITION & COMPETITIVE ANALYSIS (Market share, competitive advantages)
                6. OPERATIONAL EFFICIENCY ANALYSIS (Cost structure, productivity, scalability)
                7. RISK ASSESSMENT (Financial, operational, market, regulatory risks)
                8. VALUATION INSIGHTS (Multiples, comparables, value drivers)
                9. STRATEGIC OPPORTUNITIES (Growth drivers, expansion opportunities)
                10. KEY CONCERNS & RED FLAGS (Issues requiring attention)
                11. RECOMMENDATIONS (Strategic and tactical recommendations)
                """,
                "json_structure": """
                {
                    "executive_summary": "string",
                    "financial_performance": "string",
                    "financial_position": "string", 
                    "cash_flow_analysis": "string",
                    "market_position": "string",
                    "operational_efficiency": "string",
                    "risk_assessment": "string",
                    "valuation_insights": "string",
                    "strategic_opportunities": "string",
                    "key_concerns": "string",
                    "recommendations": ["list of recommendations"],
                    "key_metrics": {"metric_name": "value and analysis"}
                }
                """
            },
            "financial_performance": {
                "instructions": """
                Focus specifically on financial performance metrics, trends, and drivers. Analyze revenue growth, 
                profitability trends, margin analysis, and key financial ratios. This analysis should provide deep 
                insights into the company's financial performance trajectory.
                """,
                "structure": """
                1. REVENUE ANALYSIS (Growth trends, drivers, seasonality, recurring vs. one-time)
                2. PROFITABILITY ANALYSIS (Gross, operating, net margins with trends)
                3. EBITDA & OPERATING PERFORMANCE (EBITDA margins, operating leverage)
                4. FINANCIAL RATIOS ANALYSIS (Liquidity, profitability, efficiency ratios)
                5. WORKING CAPITAL ANALYSIS (Cash conversion cycle, working capital trends)
                6. COST STRUCTURE ANALYSIS (Fixed vs. variable costs, cost efficiency)
                7. PERFORMANCE BENCHMARKING (Peer comparison, industry averages)
                8. EARNINGS QUALITY ASSESSMENT (Sustainability, one-time items)
                """,
                "json_structure": """
                {
                    "revenue_analysis": "string",
                    "profitability_analysis": "string",
                    "ebitda_performance": "string",
                    "financial_ratios": "string",
                    "working_capital": "string",
                    "cost_structure": "string",
                    "benchmarking": "string",
                    "earnings_quality": "string",
                    "key_metrics": {"metric_name": "value and analysis"},
                    "recommendations": ["list of recommendations"]
                }
                """
            },
            "risk_assessment": {
                "instructions": """
                Conduct a comprehensive risk assessment covering financial, operational, market, and regulatory risks. 
                Identify key risk factors that could impact the company's performance and provide risk mitigation strategies.
                """,
                "structure": """
                1. FINANCIAL RISKS (Credit risk, liquidity risk, interest rate risk)
                2. OPERATIONAL RISKS (Business continuity, key person risk, operational failures)
                3. MARKET RISKS (Competition, market dynamics, customer concentration)
                4. REGULATORY & COMPLIANCE RISKS (Regulatory changes, compliance failures)
                5. CYBERSECURITY & TECHNOLOGY RISKS (Data breaches, system failures)
                6. ESG RISKS (Environmental, social, governance risks)
                7. SUPPLY CHAIN RISKS (Supplier dependencies, disruptions)
                8. RISK MITIGATION STRATEGIES (Risk management recommendations)
                """,
                "json_structure": """
                {
                    "financial_risks": "string",
                    "operational_risks": "string",
                    "market_risks": "string",
                    "regulatory_risks": "string",
                    "technology_risks": "string",
                    "esg_risks": "string",
                    "supply_chain_risks": "string",
                    "risk_mitigation": "string",
                    "risk_ratings": {"risk_category": "rating and explanation"},
                    "recommendations": ["list of recommendations"]
                }
                """
            },
            "market_analysis": {
                "instructions": """
                Analyze the company's market position, competitive landscape, and market dynamics. Focus on market share, 
                competitive advantages, customer analysis, and market trends that could impact the company's performance.
                """,
                "structure": """
                1. INDUSTRY OVERVIEW (Market size, growth trends, key drivers)
                2. COMPETITIVE LANDSCAPE (Key competitors, market share, positioning)
                3. COMPETITIVE ADVANTAGES (Unique value propositions, moats)
                4. CUSTOMER ANALYSIS (Customer base, retention, satisfaction)
                5. MARKET TRENDS (Emerging trends, disruptions, opportunities)
                6. PRICING ANALYSIS (Pricing power, competitive pricing)
                7. DISTRIBUTION & SALES CHANNELS (Channel effectiveness, reach)
                8. INNOVATION & R&D (Technology edge, product development)
                """,
                "json_structure": """
                {
                    "industry_overview": "string",
                    "competitive_landscape": "string",
                    "competitive_advantages": "string",
                    "customer_analysis": "string",
                    "market_trends": "string",
                    "pricing_analysis": "string",
                    "distribution_channels": "string",
                    "innovation_analysis": "string",
                    "market_metrics": {"metric_name": "value and analysis"},
                    "recommendations": ["list of recommendations"]
                }
                """
            },
            "valuation": {
                "instructions": """
                Provide comprehensive valuation analysis using multiple methodologies. Include comparable company analysis, 
                precedent transactions, and DCF considerations. Focus on valuation multiples, value drivers, and valuation insights.
                """,
                "structure": """
                1. VALUATION MULTIPLES (P/E, EV/EBITDA, P/B, EV/Revenue)
                2. COMPARABLE COMPANY ANALYSIS (Peer comparison, relative valuation)
                3. PRECEDENT TRANSACTIONS (M&A multiples, transaction comparables)
                4. DCF CONSIDERATIONS (Cash flow projections, discount rates)
                5. ASSET VALUATION (Tangible and intangible asset values)
                6. VALUE DRIVERS (Key factors driving valuation)
                7. VALUATION RANGE (Fair value range, sensitivity analysis)
                8. VALUATION RISKS (Factors that could impact valuation)
                """,
                "json_structure": """
                {
                    "valuation_multiples": "string",
                    "comparable_analysis": "string",
                    "transaction_analysis": "string",
                    "dcf_considerations": "string",
                    "asset_valuation": "string",
                    "value_drivers": "string",
                    "valuation_range": "string",
                    "valuation_risks": "string",
                    "valuation_metrics": {"metric_name": "value and analysis"},
                    "recommendations": ["list of recommendations"]
                }
                """
            }
        }
        
        return prompts.get(analysis_type, prompts["comprehensive"])
    
    def process_bulk_questions(self, questions_df: pd.DataFrame) -> pd.DataFrame:
        """Process multiple questions in bulk"""
        
        results = []
        
        for index, row in questions_df.iterrows():
            question = row.get('question', row.get('Question', ''))
            company = row.get('company', row.get('Company', ''))
            
            if not question.strip():
                continue
            
            try:
                # Search for relevant documents
                search_query = f"{company} {question}" if company else question
                search_results = self.search_financial_data(search_query, k=5)
                # Generate answer
                answer = self.generate_bulk_answer(question, search_results)
                results.append({
                    'Question': question,
                    'Company': company,
                    'Answer': answer,
                    'Sources': len(search_results),
                    'Timestamp': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                })
                
            except Exception as e:
                logger.error(f"Error processing question {index}: {str(e)}")
                results.append({
                    'Question': question,
                    'Company': company,
                    'Answer': f"Error processing question: {str(e)}",
                    'Sources': 0,
                    'Timestamp': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                })
        
        return pd.DataFrame(results)
    
    def generate_bulk_answer(self, question: str, search_results: List[Dict]) -> str:
        """Generate answer for bulk processing"""
        
        # Prepare context
        context_parts = []
        for doc in search_results[:5]:
            content = doc.get("content", "")
            file_name = doc.get("file_name", "Unknown")
            page_num = doc.get("page_number", "Unknown")
            context_parts.append(f"Source: {file_name} (Page {page_num})\n{content}")
        
        context_text = "\n\n".join(context_parts)
        
        prompt = f"""
        You are a financial analyst providing direct answers to business questions based on available documents.
        
        Question: {question}
        
        Context:
        {context_text}
        
        Provide a clear, concise answer based on the available information. If information is not available, 
        state that clearly. Focus on factual information and specific data points.
        """
        
        # Input safety gate
        if not is_text_safe(question):
            return "Question blocked by content safety policies."

        try:
            response = self.openai_client.chat.completions.create(
                model=self.deployment_name,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=1000,
                temperature=0.2
            )

            answer = response.choices[0].message.content

            # Groundedness check
            contexts = [doc.get("content", "") for doc in search_results[:6]]
            if not groundedness_check(self.openai_client, self.deployment_name, answer, contexts):
                answer = (
                    "I don't have enough evidence in the indexed documents to confidently answer this. "
                    "Please rephrase or provide more context."
                )

            # Output safety
            return enforce_output_safety(answer)

        except Exception as e:
            logger.error(f"Error generating bulk answer: {str(e)}")
            return f"Error generating answer: {str(e)}"

class PPTXGenerator:
    """Generate PowerPoint presentations for financial insights"""
    
    def create_comprehensive_presentation(self, insights_data: Dict[str, Any]) -> io.BytesIO:
        """Create a comprehensive PowerPoint presentation"""
        
        prs = Presentation()
        
        # Set slide dimensions
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        insights = insights_data.get("insights", {})
        company_name = insights_data.get("company_name", "Company")
        analysis_type = insights_data.get("analysis_type", "Analysis")
        
        # Create slides based on analysis type
        if analysis_type == "comprehensive":
            self.create_comprehensive_slides(prs, company_name, insights)
        elif analysis_type == "financial_performance":
            self.create_financial_performance_slides(prs, company_name, insights)
        elif analysis_type == "risk_assessment":
            self.create_risk_assessment_slides(prs, company_name, insights)
        elif analysis_type == "market_analysis":
            self.create_market_analysis_slides(prs, company_name, insights)
        elif analysis_type == "valuation":
            self.create_valuation_slides(prs, company_name, insights)
        
        # Save to BytesIO
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        return pptx_buffer
    
    def create_comprehensive_slides(self, prs: Presentation, company_name: str, insights: Dict):
        """Create comprehensive analysis slides"""
        
        # Title slide
        self.create_title_slide(prs, f"Financial Due Diligence Report", 
                               f"{company_name}\nComprehensive Analysis")
        
        # Executive Summary
        self.create_content_slide(prs, "Executive Summary", 
                                insights.get("executive_summary", ""))
        
        # Financial Performance
        self.create_content_slide(prs, "Financial Performance Analysis", 
                                insights.get("financial_performance", ""))
        
        # Financial Position
        self.create_content_slide(prs, "Financial Position Analysis", 
                                insights.get("financial_position", ""))
        
        # Cash Flow Analysis
        self.create_content_slide(prs, "Cash Flow Analysis", 
                                insights.get("cash_flow_analysis", ""))
        
        # Market Position
        self.create_content_slide(prs, "Market Position & Competitive Analysis", 
                                insights.get("market_position", ""))
        
        # Operational Efficiency
        self.create_content_slide(prs, "Operational Efficiency Analysis", 
                                insights.get("operational_efficiency", ""))
        
        # Risk Assessment
        self.create_content_slide(prs, "Risk Assessment", 
                                insights.get("risk_assessment", ""))
        
        # Valuation Insights
        self.create_content_slide(prs, "Valuation Insights", 
                                insights.get("valuation_insights", ""))
        
        # Strategic Opportunities
        self.create_content_slide(prs, "Strategic Opportunities", 
                                insights.get("strategic_opportunities", ""))
        
        # Key Concerns
        self.create_content_slide(prs, "Key Concerns & Red Flags", 
                                insights.get("key_concerns", ""))
        
        # Recommendations
        self.create_recommendations_slide(prs, "Recommendations", 
                                        insights.get("recommendations", []))
    
    def create_financial_performance_slides(self, prs: Presentation, company_name: str, insights: Dict):
        """Create financial performance specific slides"""
        
        self.create_title_slide(prs, f"Financial Performance Analysis", 
                               f"{company_name}\nFinancial Performance Deep Dive")
        
        sections = [
            ("Revenue Analysis", "revenue_analysis"),
            ("Profitability Analysis", "profitability_analysis"),
            ("EBITDA & Operating Performance", "ebitda_performance"),
            ("Financial Ratios Analysis", "financial_ratios"),
            ("Working Capital Analysis", "working_capital"),
            ("Cost Structure Analysis", "cost_structure"),
            ("Performance Benchmarking", "benchmarking"),
            ("Earnings Quality Assessment", "earnings_quality")
        ]
        
        for title, key in sections:
            self.create_content_slide(prs, title, insights.get(key, ""))
        
        self.create_recommendations_slide(prs, "Recommendations", 
                                        insights.get("recommendations", []))
    
    def create_risk_assessment_slides(self, prs: Presentation, company_name: str, insights: Dict):
        """Create risk assessment specific slides"""
        
        self.create_title_slide(prs, f"Risk Assessment Report", 
                               f"{company_name}\nComprehensive Risk Analysis")
        
        sections = [
            ("Financial Risks", "financial_risks"),
            ("Operational Risks", "operational_risks"),
            ("Market Risks", "market_risks"),
            ("Regulatory & Compliance Risks", "regulatory_risks"),
            ("Technology & Cybersecurity Risks", "technology_risks"),
            ("ESG Risks", "esg_risks"),
            ("Supply Chain Risks", "supply_chain_risks"),
            ("Risk Mitigation Strategies", "risk_mitigation")
        ]
        
        for title, key in sections:
            self.create_content_slide(prs, title, insights.get(key, ""))
        
        self.create_recommendations_slide(prs, "Risk Management Recommendations", 
                                        insights.get("recommendations", []))
    
    def create_market_analysis_slides(self, prs: Presentation, company_name: str, insights: Dict):
        """Create market analysis specific slides"""
        
        self.create_title_slide(prs, f"Market Analysis Report", 
                               f"{company_name}\nMarket Position & Competitive Analysis")
        
        sections = [
            ("Industry Overview", "industry_overview"),
            ("Competitive Landscape", "competitive_landscape"),
            ("Competitive Advantages", "competitive_advantages"),
            ("Customer Analysis", "customer_analysis"),
            ("Market Trends", "market_trends"),
            ("Pricing Analysis", "pricing_analysis"),
            ("Distribution & Sales Channels", "distribution_channels"),
            ("Innovation & R&D Analysis", "innovation_analysis")
        ]
        
        for title, key in sections:
            self.create_content_slide(prs, title, insights.get(key, ""))
        
        self.create_recommendations_slide(prs, "Market Strategy Recommendations", 
                                        insights.get("recommendations", []))
    
    def create_valuation_slides(self, prs: Presentation, company_name: str, insights: Dict):
        """Create valuation specific slides"""
        
        self.create_title_slide(prs, f"Valuation Analysis Report", 
                               f"{company_name}\nComprehensive Valuation Analysis")
        
        sections = [
            ("Valuation Multiples", "valuation_multiples"),
            ("Comparable Company Analysis", "comparable_analysis"),
            ("Transaction Analysis", "transaction_analysis"),
            ("DCF Considerations", "dcf_considerations"),
            ("Asset Valuation", "asset_valuation"),
            ("Value Drivers", "value_drivers"),
            ("Valuation Range", "valuation_range"),
            ("Valuation Risks", "valuation_risks")
        ]
        
        for title, key in sections:
            self.create_content_slide(prs, title, insights.get(key, ""))
        
        self.create_recommendations_slide(prs, "Valuation Recommendations", 
                                        insights.get("recommendations", []))
    
    def create_title_slide(self, prs: Presentation, title: str, subtitle: str):
        """Create title slide"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = f"{subtitle}\n\n{datetime.now().strftime('%B %Y')}"
        
        # Style title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        # Style subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    
    def create_content_slide(self, prs: Presentation, title: str, content: str):
        """Create content slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        # Add content
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Split content into bullet points
        points = content.split('. ')
        for point in points:
            if point.strip():
                p = text_frame.add_paragraph()
                p.text = point.strip()
                p.font.size = Pt(16)
                p.space_after = Pt(12)
    
    def create_recommendations_slide(self, prs: Presentation, title: str, recommendations: List[str]):
        """Create recommendations slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        # Add recommendations
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, recommendation in enumerate(recommendations, 1):
            p = text_frame.add_paragraph()
            p.text = f"{i}. {recommendation}"
            p.font.size = Pt(16)
            p.space_after = Pt(12)

class BulkDocumentGenerator:
    """Generate bulk documents for Q&A"""
    
    def create_word_document(self, df: pd.DataFrame) -> io.BytesIO:
        """Create Word document from DataFrame"""
        
        doc = Document()
        
        # Add title
        title = doc.add_heading('Financial Due Diligence Q&A Report', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add timestamp
        doc.add_paragraph(f"Generated on: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}")
        doc.add_paragraph("")
        
        # Add Q&A content
        for index, row in df.iterrows():
            # Question
            question_para = doc.add_paragraph()
            question_run = question_para.add_run(f"Q{index+1}: {row['Question']}")
            question_run.bold = True
            question_run.font.size = Pt(12)
            
            # Company (if available)
            if row.get('Company') and row['Company'].strip():
                company_para = doc.add_paragraph()
                company_run = company_para.add_run(f"Company: {row['Company']}")
                company_run.italic = True
            
            # Answer
            answer_para = doc.add_paragraph()
            answer_run = answer_para.add_run(f"Answer: {row['Answer']}")
            answer_run.font.size = Pt(11)
            
            # Sources
            sources_para = doc.add_paragraph()
            sources_run = sources_para.add_run(f"Sources: {row['Sources']} documents")
            sources_run.font.size = Pt(10)
            sources_run.italic = True
            
            # Add spacing
            doc.add_paragraph("")
        
        # Save to BytesIO
        doc_buffer = io.BytesIO()
        doc.save(doc_buffer)
        doc_buffer.seek(0)
        
        return doc_buffer

def main():
    st.title("💼 Financial Due Diligence Insights")
    st.markdown("### Professional Analysis for Transaction Services & Due Diligence Teams")
    
    # Check if index exists
    index_name = st.session_state.get('index_name')
    if not index_name:
        st.warning("⚠️ Please run the data pipeline first to create a search index.")
        return
    
    # Initialize clients
    try:
        #credential = DefaultAzureCredential()
        credential = get_search_credential()
        # openai_client = AzureOpenAI(
        #     azure_ad_token_provider=lambda: credential.get_token("https://cognitiveservices.azure.com/.default").token,
        #     azure_endpoint=OPENAI_ENDPOINT_URL,
        #     api_version="2024-02-01"
        # )
        openai_client = get_openai_client()
        
        search_client = SearchClient(
            endpoint=SEARCH_SERVICE_ENDPOINT,
            index_name=index_name,
            credential=credential
        )
        
        # Initialize agents
        insights_agent = FinancialInsightsAgent(search_client, openai_client, DEPLOYMENT_NAME)
        pptx_generator = PPTXGenerator()
        bulk_generator = BulkDocumentGenerator()
        
    except Exception as e:
        st.error(f"❌ Error initializing clients: {str(e)}")
        return
    
    # Tab selection
    tab1, tab2 = st.tabs(["📊 Single Company Analysis", "📋 Bulk Question Processing"])
    
    with tab1:
        # Single company analysis
        st.markdown("### 📊 Comprehensive Financial Analysis")
        
        col1, col2 = st.columns(2)
        
        with col1:
            company_name = st.text_input(
                "Company Name:",
                placeholder="e.g., ABC Corporation",
                help="Enter the company name for analysis"
            )
            
            analysis_type = st.selectbox(
                "Analysis Type:",
                options=[
                    "comprehensive",
                    "financial_performance", 
                    "risk_assessment",
                    "market_analysis",
                    "valuation"
                ],
                format_func=lambda x: {
                    "comprehensive": "📊 Comprehensive Due Diligence Agent",
                    "financial_performance": "💰 Financial Performance Deep Dive Agent",
                    "risk_assessment": "⚠️ Risk Assessment & Mitigation Agent",
                    "market_analysis": "🎯 Market Position & Competitive Analysis Agent",
                    "valuation": "💎 Valuation Analysis & Multiples Agent"
                }[x],
                help="Select the type of financial analysis"
            )
        
        with col2:
            custom_focus = st.text_area(
                "Custom Focus Areas:",
                placeholder="e.g., ESG factors, digital transformation, supply chain risks",
                help="Add specific areas of focus for the analysis"
            )
        
        # Generate insights
        if st.button("🚀 Generate Comprehensive Analysis", type="primary"):
            if not company_name.strip():
                st.error("Please enter a company name.")
                return
            
            with st.spinner("🔍 Analyzing financial data and generating comprehensive insights..."):
                try:
                    # Generate insights
                    insights_data = insights_agent.generate_comprehensive_insights(
                        company_name, analysis_type, custom_focus
                    )
                    
                    # Store in session state
                    st.session_state['financial_insights'] = insights_data
                    
                    # Display comprehensive insights
                    st.markdown("### 📈 Comprehensive Financial Analysis")
                    st.success(f"✅ {analysis_type.replace('_', ' ').title()} analysis completed for {company_name}")
                    
                    insights = insights_data.get("insights", {})
                    
                    # Display all sections based on analysis type
                    if analysis_type == "comprehensive":
                        sections = [
                            ("📋 Executive Summary", "executive_summary"),
                            ("💰 Financial Performance", "financial_performance"),
                            ("🏦 Financial Position", "financial_position"),
                            ("💸 Cash Flow Analysis", "cash_flow_analysis"),
                            ("🎯 Market Position", "market_position"),
                            ("⚙️ Operational Efficiency", "operational_efficiency"),
                            ("⚠️ Risk Assessment", "risk_assessment"),
                            ("💎 Valuation Insights", "valuation_insights"),
                            ("🚀 Strategic Opportunities", "strategic_opportunities"),
                            ("🚨 Key Concerns", "key_concerns")
                        ]
                    elif analysis_type == "financial_performance":
                        sections = [
                            ("📊 Revenue Analysis", "revenue_analysis"),
                            ("💰 Profitability Analysis", "profitability_analysis"),
                            ("📈 EBITDA Performance", "ebitda_performance"),
                            ("🔢 Financial Ratios", "financial_ratios"),
                            ("🔄 Working Capital", "working_capital"),
                            ("💸 Cost Structure", "cost_structure"),
                            ("📊 Benchmarking", "benchmarking"),
                            ("✅ Earnings Quality", "earnings_quality")
                        ]
                    elif analysis_type == "risk_assessment":
                        sections = [
                            ("💰 Financial Risks", "financial_risks"),
                            ("⚙️ Operational Risks", "operational_risks"),
                            ("🎯 Market Risks", "market_risks"),
                            ("⚖️ Regulatory Risks", "regulatory_risks"),
                            ("🔒 Technology Risks", "technology_risks"),
                            ("🌱 ESG Risks", "esg_risks"),
                            ("🔗 Supply Chain Risks", "supply_chain_risks"),
                            ("🛡️ Risk Mitigation", "risk_mitigation")
                        ]
                    elif analysis_type == "market_analysis":
                        sections = [
                            ("🏭 Industry Overview", "industry_overview"),
                            ("🎯 Competitive Landscape", "competitive_landscape"),
                            ("💪 Competitive Advantages", "competitive_advantages"),
                            ("👥 Customer Analysis", "customer_analysis"),
                            ("📈 Market Trends", "market_trends"),
                            ("💰 Pricing Analysis", "pricing_analysis"),
                            ("🚚 Distribution Channels", "distribution_channels"),
                            ("💡 Innovation Analysis", "innovation_analysis")
                        ]
                    elif analysis_type == "valuation":
                        sections = [
                            ("📊 Valuation Multiples", "valuation_multiples"),
                            ("🏢 Comparable Analysis", "comparable_analysis"),
                            ("📈 Transaction Analysis", "transaction_analysis"),
                            ("💰 DCF Considerations", "dcf_considerations"),
                            ("🏗️ Asset Valuation", "asset_valuation"),
                            ("🎯 Value Drivers", "value_drivers"),
                            ("📊 Valuation Range", "valuation_range"),
                            ("⚠️ Valuation Risks", "valuation_risks")
                        ]
                    
                    # Display all sections
                    for section_title, section_key in sections:
                        if section_key in insights and insights[section_key]:
                            st.markdown(f"#### {section_title}")
                            st.markdown(insights[section_key])
                            st.markdown("---")
                    
                    # Key Metrics
                    metrics_key = "key_metrics" if analysis_type == "comprehensive" else f"{analysis_type.split('_')[0]}_metrics"
                    if metrics_key in insights and insights[metrics_key]:
                        st.markdown("#### 📊 Key Metrics")
                        metrics_col1, metrics_col2 = st.columns(2)
                        
                        metrics_list = list(insights[metrics_key].items())
                        mid_point = len(metrics_list) // 2
                        
                        with metrics_col1:
                            for metric, value in metrics_list[:mid_point]:
                                st.metric(metric, value)
                        
                        with metrics_col2:
                            for metric, value in metrics_list[mid_point:]:
                                st.metric(metric, value)
                    
                    # Recommendations
                    if "recommendations" in insights and insights["recommendations"]:
                        st.markdown("#### 💡 Recommendations")
                        for i, rec in enumerate(insights["recommendations"], 1):
                            st.markdown(f"{i}. {rec}")
                    
                except Exception as e:
                    st.error(f"❌ Error generating insights: {str(e)}")
                    logger.error(f"Error generating insights: {str(e)}")
        
        # Download section
        if 'financial_insights' in st.session_state:
            st.markdown("---")
            st.markdown("### 📥 Download Professional Report")
            
            if st.button("📊 Download PowerPoint Presentation (PPTX)", type="secondary"):
                with st.spinner("Creating comprehensive PowerPoint presentation..."):
                    try:
                        pptx_buffer = pptx_generator.create_comprehensive_presentation(
                            st.session_state['financial_insights']
                        )
                        
                        # Encode for download
                        pptx_b64 = base64.b64encode(pptx_buffer.read()).decode()
                        
                        company_name = st.session_state['financial_insights'].get('company_name', 'Company')
                        analysis_type = st.session_state['financial_insights'].get('analysis_type', 'analysis')
                        filename = f"{company_name}_{analysis_type}_Analysis_{datetime.now().strftime('%Y%m%d')}.pptx"
                        
                        st.markdown(
                            f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{pptx_b64}" download="{filename}">📊 Download {filename}</a>',
                            unsafe_allow_html=True
                        )
                        
                        st.success("✅ Professional PowerPoint presentation ready for download!")
                        
                    except Exception as e:
                        st.error(f"❌ Error creating PowerPoint: {str(e)}")
    
    with tab2:
        # Bulk processing
        st.markdown("### 📋 Bulk Question Processing")
        st.markdown("Upload an Excel file with questions and get comprehensive answers")
        
        # File upload
        uploaded_file = st.file_uploader(
            "Upload Excel file with questions",
            type=['xlsx', 'xls'],
            help="Excel should have columns: 'Question' (required), 'Company' (optional)"
        )
        
        if uploaded_file is not None:
            try:
                # Read Excel file
                df = pd.read_excel(uploaded_file)
                
                # Display preview
                st.markdown("#### 📊 File Preview")
                st.dataframe(df.head())
                
                # Validate columns
                if 'Question' not in df.columns and 'question' not in df.columns:
                    st.error("❌ Excel file must contain a 'Question' column")
                    return
                
                st.success(f"✅ Found {len(df)} questions to process")
                
                # Process questions
                if st.button("🚀 Process All Questions", type="primary"):
                    with st.spinner("Processing questions... This may take a few minutes."):
                        try:
                            # Process bulk questions
                            results_df = insights_agent.process_bulk_questions(df)
                            
                            # Store results
                            st.session_state['bulk_results'] = results_df
                            
                            # Display results
                            st.markdown("#### 📊 Processing Results")
                            st.dataframe(results_df)
                            
                            # Download options
                            st.markdown("#### 📥 Download Options")
                            
                            col1, col2 = st.columns(2)
                            
                            with col1:
                                # Excel download
                                excel_buffer = io.BytesIO()
                                results_df.to_excel(excel_buffer, index=False)
                                excel_buffer.seek(0)
                                
                                st.download_button(
                                    label="📊 Download Excel Report",
                                    data=excel_buffer,
                                    file_name=f"bulk_qa_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                                )
                            
                            with col2:
                                # Word download
                                word_buffer = bulk_generator.create_word_document(results_df)
                                
                                st.download_button(
                                    label="📄 Download Word Report",
                                    data=word_buffer,
                                    file_name=f"bulk_qa_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx",
                                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                                )
                            
                        except Exception as e:
                            st.error(f"❌ Error processing questions: {str(e)}")
                            logger.error(f"Error processing bulk questions: {str(e)}")
                            
            except Exception as e:
                st.error(f"❌ Error reading Excel file: {str(e)}")
        
        # Template download
        st.markdown("---")
        st.markdown("#### 📥 Download Template")
        
        # Create template
        template_data = {
            'Question': [
                'What are the key financial metrics for Q3 2024?',
                'What are the main risk factors identified?',
                'What is the competitive positioning?'
            ],
            'Company': [
                'ABC Corporation',
                'ABC Corporation',
                'ABC Corporation'
            ]
        }
        
        template_df = pd.DataFrame(template_data)
        template_buffer = io.BytesIO()
        template_df.to_excel(template_buffer, index=False)
        template_buffer.seek(0)
        
        st.download_button(
            label="📊 Download Excel Template",
            data=template_buffer,
            file_name="bulk_questions_template.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

if __name__ == "__main__":
    main()
