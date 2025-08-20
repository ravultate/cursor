import base64
import os
import time
import json
import logging
from typing import List, Dict
from io import BytesIO
import requests
import fitz  # PyMuPDF
import streamlit as st
from dotenv import load_dotenv
from docx import Document
from pptx import Presentation
import pandas as pd
from azure.core.credentials import AzureKeyCredential
from azure.core.exceptions import ResourceNotFoundError
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient
from azure.search.documents.indexes.models import (
    SimpleField, SearchFieldDataType, SearchableField, SearchField,
    VectorSearch, HnswAlgorithmConfiguration, VectorSearchProfile,
    SemanticConfiguration, SemanticPrioritizedFields, SemanticField,
    SemanticSearch, SearchIndex, AzureOpenAIVectorizer,
    AzureOpenAIVectorizerParameters, ScoringProfile, TextWeights
)
from openai import AzureOpenAI

from auth_utils import get_search_credential,get_openai_client

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("agentic_rag.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

load_dotenv()

SEARCH_SERVICE_ENDPOINT = os.getenv("SEARCH_SERVICE_ENDPOINT")
OPENAI_ENDPOINT_URL = os.getenv("OPENAI_ENDPOINT_URL")
EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
SEARCH_API_KEY = os.getenv("SEARCH_API_KEY")

# ──────────────── SharePoint Files and Folders Browsing Utils ────────────────

def get_access_token():
    token = st.session_state.get("access_token")
    if not token:
        st.error("Please login first to get access token.")
        st.stop()
    return token

def list_drive_children(drive_id: str, folder_id: str, access_token: str):
    """
    Returns files and folders in the given SharePoint folder.
    If folder_id is None, lists root.
    """
    headers = {"Authorization": f"Bearer {access_token}"}
    if folder_id:
        url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{folder_id}/children"
    else:
        url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/children"
    
    try:
        resp = requests.get(url, headers=headers, timeout=30)
        resp.raise_for_status()
        return resp.json().get("value", [])
    except requests.exceptions.RequestException as e:
        logger.error(f"Error listing drive children: {e}")
        st.error(f"Error accessing SharePoint folder: {e}")
        return []

def get_documents_drive_id(site_id: str, access_token: str):
    headers = {"Authorization": f"Bearer {access_token}"}
    url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives"
    try:
        drives = requests.get(url, headers=headers, timeout=30).json().get("value", [])
        doc_drive = next((d for d in drives if d.get("name") == "Documents"), None)
        if not doc_drive:
            st.error("Documents library not found.")
            st.stop()
        return doc_drive["id"]
    except requests.exceptions.RequestException as e:
        logger.error(f"Error getting documents drive ID: {e}")
        st.error(f"Error accessing SharePoint site: {e}")
        st.stop()

def download_pdf(drive_id: str, file_id: str, access_token: str) -> bytes:
    headers = {"Authorization": f"Bearer {access_token}"}
    url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{file_id}/content"
    
    for attempt in range(3):
        try:
            resp = requests.get(url, headers=headers, timeout=60)
            resp.raise_for_status()
            return resp.content
        except requests.exceptions.RequestException as e:
            logger.warning("Download attempt %s failed: %s", attempt + 1, e)
            if attempt == 2:
                raise
            time.sleep(2 ** attempt)

def download_file_content(drive_id: str, file_id: str, access_token: str) -> bytes:
    """Generic download for any file type from SharePoint."""
    return download_pdf(drive_id, file_id, access_token)

def extract_pdf_text(pdf_bytes: bytes) -> List[tuple]:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    pages_text = []
    for page_num in range(doc.page_count):
        try:
            page = doc.load_page(page_num)
            text = page.get_text("text")
            if text.strip():
                pages_text.append((page_num + 1, text))
        except Exception as e:
            logger.warning("Error extracting text from page %s: %s", page_num + 1, e)
    doc.close()
    return pages_text

def extract_docx_text(docx_bytes: bytes) -> List[tuple]:
    try:
        document = Document(BytesIO(docx_bytes))
        text_lines = [p.text.strip() for p in document.paragraphs if p.text and p.text.strip()]
        joined = "\n".join(text_lines)
        return [(1, joined)] if joined.strip() else []
    except Exception as e:
        logger.warning("Error extracting DOCX text: %s", e)
        return []

def extract_pptx_text(pptx_bytes: bytes) -> List[tuple]:
    try:
        prs = Presentation(BytesIO(pptx_bytes))
        pages = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in paragraph.runs).strip()
                        if line:
                            texts.append(line)
            slide_text = "\n".join(texts)
            if slide_text.strip():
                pages.append((idx, slide_text))
        return pages
    except Exception as e:
        logger.warning("Error extracting PPTX text: %s", e)
        return []

def extract_csv_text(csv_bytes: bytes) -> List[tuple]:
    try:
        # Let pandas detect encoding; fallback handled by pandas
        df = pd.read_csv(BytesIO(csv_bytes))
        df = df.fillna("")
        text = df.to_csv(index=False)
        return [(1, text)] if text.strip() else []
    except Exception as e:
        logger.warning("Error extracting CSV text: %s", e)
        return []

def extract_excel_text(excel_bytes: bytes) -> List[tuple]:
    try:
        xls = pd.ExcelFile(BytesIO(excel_bytes))
        pages: List[tuple] = []
        for i, sheet in enumerate(xls.sheet_names, start=1):
            df = xls.parse(sheet)
            df = df.fillna("")
            text = f"Sheet: {sheet}\n" + df.to_csv(index=False)
            if text.strip():
                pages.append((i, text))
        return pages
    except Exception as e:
        logger.warning("Error extracting Excel text: %s", e)
        return []

def extract_text_by_extension(file_name: str, file_bytes: bytes) -> List[tuple]:
    ext = os.path.splitext(file_name)[1].lower()
    if ext == ".pdf":
        return extract_pdf_text(file_bytes)
    if ext == ".docx":
        return extract_docx_text(file_bytes)
    if ext == ".pptx":
        return extract_pptx_text(file_bytes)
    if ext in (".xlsx", ".xls"):
        return extract_excel_text(file_bytes)
    if ext == ".csv":
        return extract_csv_text(file_bytes)
    return []

def intelligent_chunk_text(text: str, max_chunk_size: int = 1000, overlap: int = 200):
    if not text.strip():
        return []
    
    sentences = text.replace("\n", " ").split(". ")
    chunks = []
    current_chunk = ""
    
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        
        if len(current_chunk) + len(sentence) + 2 > max_chunk_size and current_chunk:
            chunks.append(current_chunk.strip())
            current_chunk = (current_chunk[-overlap:] + " " + sentence) if overlap and len(current_chunk) > overlap else sentence
        else:
            current_chunk = (". " if current_chunk else "") + sentence
    
    if current_chunk:
        chunks.append(current_chunk.strip())
    
    return chunks

def generate_embeddings_batch(openai_client: AzureOpenAI, texts: List[str], batch_size: int = 16):
    embeddings = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        try:
            resp = openai_client.embeddings.create(input=batch, model=EMBEDDING_MODEL)
            embeddings.extend([item.embedding for item in resp.data])
            time.sleep(0.1)
        except Exception as e:
            logger.error("Error generating embeddings batch %s: %s", i // batch_size + 1, e)
            embeddings.extend([[0.0] * 3072] * len(batch))
    return embeddings

def create_enhanced_search_index(index_name: str):
    #credential = DefaultAzureCredential()
    credential = get_search_credential()
    index_client = SearchIndexClient(endpoint=SEARCH_SERVICE_ENDPOINT, credential=credential)
    
    fields = [
        SimpleField(name="id", type=SearchFieldDataType.String, key=True, sortable=True, filterable=True),
        SearchableField(name="content", type=SearchFieldDataType.String, searchable=True, retrievable=True, analyzer_name="standard.lucene"),
        SearchField(name="content_vector", type=SearchFieldDataType.Collection(SearchFieldDataType.Single),
                   searchable=True, stored=True, vector_search_dimensions=3072, vector_search_profile_name="myHnswProfile"),
        SimpleField(name="page_number", type=SearchFieldDataType.Int32, filterable=True, sortable=True),
        SearchField(name="file_name", type="Edm.String", retrievable=True),
        SimpleField(name="chunk_id", type=SearchFieldDataType.String, filterable=True, retrievable=True),
        SimpleField(name="document_type", type=SearchFieldDataType.String, filterable=True),
        SimpleField(name="last_modified", type=SearchFieldDataType.DateTimeOffset, filterable=True, sortable=True)
    ]
    
    vector_search = VectorSearch(
        algorithms=[HnswAlgorithmConfiguration(name="myHnsw",
                                             parameters={"m": 4, "ef_construction": 400, "ef_search": 500, "metric": "cosine"})],
        profiles=[VectorSearchProfile(name="myHnswProfile", algorithm_configuration_name="myHnsw", vectorizer_name="myVectorizer")],
        vectorizers=[AzureOpenAIVectorizer(
            vectorizer_name="myVectorizer",
            parameters=AzureOpenAIVectorizerParameters(
                resource_url=OPENAI_ENDPOINT_URL,
                deployment_name=EMBEDDING_MODEL,
                model_name=EMBEDDING_MODEL,
            ),
        )],
    )
    
    semantic_config = SemanticConfiguration(
        name="semantic_config",
        prioritized_fields=SemanticPrioritizedFields(
            title_field=SemanticField(field_name="file_name"),
            content_fields=[SemanticField(field_name="content")],
            keywords_fields=[SemanticField(field_name="document_type")],
        ),
    )
    
    semantic_search = SemanticSearch(configurations=[semantic_config])
    
    scoring_profile = ScoringProfile(
        name="freshness_boost",
        text_weights=TextWeights(weights={"content": 1.0, "file_name": 2.0}),
        functions=[{
            "type": "freshness",
            "fieldName": "last_modified",
            "boost": 2.0,
            "interpolation": "linear",
            "freshness": {"boostingDuration": "P30D"}
        }],
    )
    
    index = SearchIndex(
        name=index_name,
        fields=fields,
        vector_search=vector_search,
        semantic_search=semantic_search,
        scoring_profiles=[scoring_profile],
    )
    
    index_client.create_or_update_index(index)

def upload_to_search_index_batch(index_name: str, docs: List[Dict], credential:None):
    if credential is None:
        credential = get_search_credential()

    client = SearchClient(endpoint=SEARCH_SERVICE_ENDPOINT, index_name=index_name, credential=credential)
    batch_size = 1000
    total = len(docs)
    uploaded = 0
    
    for i in range(0, total, batch_size):
        batch = docs[i:i + batch_size]
        try:
            results = client.upload_documents(documents=batch)
            uploaded = sum(1 for r in results if r.succeeded)
            failed = [r for r in results if not r.succeeded]
            if failed:
                logger.warning("Batch %s: %s docs failed", i // batch_size + 1, len(failed))
        except Exception as e:
            logger.error("Upload batch starting %s failed: %s", i, e)
        time.sleep(0.1)
    
    return uploaded == total

def get_all_selected_files_from_all_folders():
    """
    Helper function to collect all selected files from all folders
    """
    all_selected_files = []
    file_selections = st.session_state.get("file_selections", {})
    folder_file_mapping = st.session_state.get("folder_file_mapping", {})
    
    for folder_id, selected_file_names in file_selections.items():
        if selected_file_names and folder_id in folder_file_mapping:
            folder_files = folder_file_mapping[folder_id]
            selected_file_objs = [file for file in folder_files if file["name"] in selected_file_names]
            all_selected_files.extend(selected_file_objs)
    
    return all_selected_files

def render_navigation_bar(current_folder_name: str, has_parent: bool):
    """Render the navigation bar with consistent styling"""
    st.markdown('<div class="nav-container">', unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([3, 1, 1])
    
    with col1:
        # Enhanced breadcrumb with better styling
        path_parts = [f[1] for f in st.session_state.folder_stack] or ["📁 Root"]
        path_display = " → ".join(path_parts)
        st.markdown(f'<div class="breadcrumb-container"><strong>📍 Current Location:</strong> {path_display}</div>', unsafe_allow_html=True)
    
    with col1:
        if has_parent:
            if st.button("⬅️ Go Back", help="Navigate to parent folder", key="nav_back"):
                st.session_state.folder_stack.pop()
                st.rerun()
    
    with col3:
        # Show total selected files across all folders
        all_selected = get_all_selected_files_from_all_folders()
        if all_selected:
            st.metric("📄 Selected", len(all_selected))
    
    st.markdown('</div>', unsafe_allow_html=True)

def render_empty_folder_message(current_folder_name: str):
    """Render a user-friendly message for empty folders with clear navigation options"""
    st.markdown("""
    <div class="empty-folder-container">
        <h3>📂 Empty Folder</h3>
        <p>No supported files found in this folder: <strong>{}</strong></p>
        <p>💡 <em>You can navigate back to explore other folders or go up to the parent directory.</em></p>
    </div>
    """.format(current_folder_name), unsafe_allow_html=True)
    
    # Prominent navigation options for empty folders
    col1, col2, col3 = st.columns([1, 2, 1])
    
    with col1:
        if st.session_state.folder_stack:
            parent_folder = st.session_state.folder_stack[-2][1] if len(st.session_state.folder_stack) > 1 else "Root"
            if st.button(f"⬅️ Back to {parent_folder}", type="primary", use_container_width=True, key="empty_folder_back"):
                st.session_state.folder_stack.pop()
                st.rerun()
        else:
            st.info("You are in the root directory. Navigate to subfolders to find supported files.")

# ──────────── ENHANCED FOLDER & FILE BROWSING LOGIC ────────────

def browse_sharepoint_folders_and_files(site_id: str, access_token: str):
    """
    Enhanced UI for folder navigation with production-ready styling and better UX for empty folders.
    """
    # Apply custom CSS
  
    
    # Setup persistent navigation state
    if "folder_stack" not in st.session_state:
        st.session_state.folder_stack = []  # Each element: (folder_id, "Folder Name")
    if "file_selections" not in st.session_state:
        st.session_state.file_selections = {}  # folder_id -> [selected_file_names]
    if "folder_file_mapping" not in st.session_state:
        st.session_state.folder_file_mapping = {}  # folder_id -> [file_objects]
    
    # Get drive_id
    drive_id = get_documents_drive_id(site_id, access_token)
    
    # Determine current folder (None is root)
    if st.session_state.folder_stack:
        current_folder_id, current_folder_name = st.session_state.folder_stack[-1]
    else:
        current_folder_id = None
        current_folder_name = "Root"
    
    # Always render navigation bar first
    render_navigation_bar(current_folder_name, bool(st.session_state.folder_stack))
    
    # List children (files and folders)
    with st.spinner(f"🔍 Loading contents of: **{current_folder_name}**"):
        children = list_drive_children(drive_id, current_folder_id, access_token)
    
    if not children:
        st.warning("⚠️ No items found in this folder or error accessing the folder.")
        render_empty_folder_message(current_folder_name)
        return [], drive_id
    
    # Separate folders and supported files
    folders = [item for item in children if item.get("folder") is not None]
    allowed_mime_types = {
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # docx
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # pptx
        "text/csv",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # xlsx
        "application/vnd.ms-excel",  # xls
    }
    files = [
        item for item in children
        if item.get("file") is not None and item["file"].get("mimeType", "") in allowed_mime_types
    ]
    
    # Store file mapping for this folder
    st.session_state.folder_file_mapping[current_folder_id] = files
    
    # Display folders section with enhanced styling
    if folders:
        st.markdown('<h3 class="section-header">📁 Folders ({} found)</h3>'.format(len(folders)), unsafe_allow_html=True)
        
        # Create columns for folder display (3 folders per row)
        folder_cols = st.columns(3)
        for idx, folder in enumerate(folders):
            col_idx = idx % 3
            with folder_cols[col_idx]:
                folder_name = folder["name"]
                # Truncate long folder names for display
                display_name = folder_name if len(folder_name) <= 20 else folder_name[:17] +  "..."
                
                if st.button(
                    f"📂 {display_name}", 
                    key=f"folder_{folder['id']}", 
                    help=f"Navigate to: {folder_name}",
                    use_container_width=True
                ):
                    st.session_state.folder_stack.append((folder["id"], folder_name))
                    st.rerun()
        
        st.markdown("---")
    
    # Display files section
    st.markdown('<h3 class="section-header">📄 Files ({} found)</h3>'.format(len(files)), unsafe_allow_html=True)
    
    if files:
        # Get current selections for this folder
        current_selections = st.session_state.file_selections.get(current_folder_id, [])
        
        # File selection UI



        col1, col2, col3 = st.columns([3, 1, 1])

        with col1:
            file_names = [file["name"] for file in files]
            selected_files = st.multiselect(
                "Select files to process:",
                options=file_names,
                default=current_selections,
                key=f"file_select_{current_folder_id}",
                help="Select one or more files to include in processing"
            )
            # Update session state
            st.session_state.file_selections[current_folder_id] = selected_files

        with col2:
            if st.button("✅ Select All", help="Select all supported files in this folder", key=f"select_all_{current_folder_id}", use_container_width=True):
                st.session_state.file_selections[current_folder_id] = file_names
                st.rerun()
            if st.button("❌ Clear All", help="Clear all selections in this folder", key=f"clear_all_{current_folder_id}", use_container_width=True):
                st.session_state.file_selections[current_folder_id] = []
                st.rerun()
            if st.session_state.folder_stack:
                parent_folder = st.session_state.folder_stack[-2][1] if len(st.session_state.folder_stack) > 1 else "Root"
                if st.button(f"⬅️ Back to {parent_folder}", type="primary", key="aligned_folder_back", use_container_width=True):
                    st.session_state.folder_stack.pop()
                    st.rerun()
        
        # Show selected files in current folder
        if selected_files:
            st.success(f"✅ {len(selected_files)} file(s) selected in current folder")
            
            # Show selected file names
            with st.expander("📋 View Selected Files in Current Folder"):
                for file_name in selected_files:
                    st.write(f"• {file_name}")
    else:
        # Enhanced empty folder handling
        render_empty_folder_message(current_folder_name)
    
    # Summary section showing all selections across folders
    all_selected_files = get_all_selected_files_from_all_folders()
    
    if all_selected_files:
        st.markdown("---")
        st.markdown('<h3 class="section-header">🎯 Selection Summary</h3>', unsafe_allow_html=True)
        
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Total Selected Files", len(all_selected_files))
        with col2:
            folders_with_selections = len([k for k, v in st.session_state.file_selections.items() if v])
            st.metric("Folders with Selections", folders_with_selections)
        with col3:
            total_folders_visited = len(st.session_state.folder_file_mapping)
            st.metric("Folders Explored", total_folders_visited)
        
        # Show detailed breakdown
        with st.expander("📋 View All Selected Files by Folder", expanded=False):
            for folder_id, selected_file_names in st.session_state.file_selections.items():
                if selected_file_names:
                    # Get folder name
                    if folder_id is None:
                        folder_name = "Root"
                    else:
                        # Find folder name from the stack or use ID
                        folder_name = next(
                            (name for fid, name in st.session_state.folder_stack if fid == folder_id),
                            f"Folder ID: {folder_id}"
                        )
                    
                    st.write(f"**📁 {folder_name}:**")
                    for file_name in selected_file_names:
                        st.write(f"  • {file_name}")
                    st.write("")
    
    return all_selected_files, drive_id

# ───── MAIN APP LOGIC ─────

def main():
    # Apply custom CSS at the start
    
    
    access_token = get_access_token()
    site = st.session_state.get("selected_site")
    
    if not site:
        st.error("Please login and select a SharePoint site first.")
        st.stop()
    
    site_id = site["id"]
    site_url = site["webUrl"]
    
    st.title("📚 SharePoint Document Processor")
    st.markdown(f"**Selected SharePoint Site:** `{site_url}`")
    
    # File browsing and selection
    pdf_files, drive_id = browse_sharepoint_folders_and_files(site_id, access_token)

    if not pdf_files:
        st.markdown("---")
        st.info("💡 **How to get started:**")
        st.markdown("""
        1. 📁 **Navigate folders:** Click on folder buttons to explore directories
        2. 📄 **Select files:** Use checkboxes to select supported files from any folder
        3. ✅ **Quick select:** Use 'Select All' or 'Clear All' for convenience
        4. 🚀 **Process:** Click 'Start Processing' when you have files selected
        """)
        st.stop()

    # Improved UI layout for processing section
    st.markdown("---")
    st.markdown('<h3 class="section-header" style="margin-bottom:0.5rem;">🚀 Ready to Process</h3>', unsafe_allow_html=True)
    st.markdown("<div style='height:10px;'></div>", unsafe_allow_html=True)

    # Metrics row
    metrics_row = st.columns([1, 1, 1])
    metrics_row[0].metric("Files Selected", len(pdf_files))
    metrics_row[1].metric("Folders Explored", len(st.session_state.folder_file_mapping))
    metrics_row[2].metric("Selections", sum(len(v) for v in st.session_state.file_selections.values()))

    st.markdown("<div style='height:10px;'></div>", unsafe_allow_html=True)

    # Action buttons row
    action_row = st.columns([2, 1])
    with action_row[0]:
        st.success(f"✅ {len(pdf_files)} file(s) selected and ready for processing")
    with action_row[1]:
        start_processing = st.button(
            "🚀 Start Processing All Selected Files",
            type="primary",
            use_container_width=True,
            help="Begin processing all selected files from all folders"
        )

    st.markdown("<div style='height:10px;'></div>", unsafe_allow_html=True)

    if not start_processing:
        st.info("👆 Click the button above to start processing the selected files.")
        return
    
    # ────────── INDEX EXISTENCE CHECK ──────────
    site_name = site_url.rstrip("/").split("/")[-1]
    index_name = f"{site_name.lower()}-enhanced-rag-index"
    
    # credential = AzureKeyCredential(SEARCH_API_KEY)
    # index_client = SearchIndexClient(endpoint=SEARCH_SERVICE_ENDPOINT, credential=credential)
    #credential = DefaultAzureCredential()
    credential = get_search_credential()

    index_client = SearchIndexClient(endpoint=SEARCH_SERVICE_ENDPOINT,
                                    credential=credential)
    try:
        index_client.get_index(index_name)
        index_exists = True
        st.info(f"ℹ️ Index **{index_name}** already exists — new documents will be added.")
    except ResourceNotFoundError:
        index_exists = False
    
    if not index_exists:
        st.info(f"🏗️ Creating Azure AI Search index: **{index_name}**")
        with st.spinner("Building index with hybrid search and semantic ranking..."):
            try:
                create_enhanced_search_index(index_name)
                st.success("✅ Index created successfully!")
            except Exception as e:
                st.error(f"❌ Error creating index: {e}")
                logger.error(f"Index creation failed: {e}")
                st.stop()
    
    # ────────── BEGIN PROCESSING ──────────
    #credential = DefaultAzureCredential()
    credential = get_search_credential()

    # openai_client = AzureOpenAI(
    #   azure_ad_token_provider=lambda: credential
    #        .get_token("https://cognitiveservices.azure.com/.default").token,
    #      azure_endpoint=OPENAI_ENDPOINT_URL,
    #      api_version="2024-02-01"
    #  )
    
    openai_client = get_openai_client()
    docs_to_upload = []
    total_files = len(pdf_files)
    
    # Progress tracking with enhanced UI
    progress_bar = st.progress(0)
    status_container = st.container()
    
    with status_container:
        st.markdown("### 📊 Processing Progress")
        metrics_cols = st.columns(4)
        
        processed_files_metric = metrics_cols[0].empty()
        chunks_generated_metric = metrics_cols[1].empty()
        current_file_metric = metrics_cols[2].empty()
        success_rate_metric = metrics_cols[3].empty()
    
    successful_files = 0
    
    for i, file_meta in enumerate(pdf_files, start=1):
        fname = file_meta["name"]
        fid = file_meta["id"]
        
        # Update current file metric
        current_file_metric.metric("📄 Current File", f"{i}/{total_files}")
        
        try:
            # Download and process file
            with st.spinner(f"⬇️ Downloading {fname}..."):
                file_bytes = download_file_content(drive_id, fid, access_token)
            
            with st.spinner(f"📖 Extracting text from {fname}..."):
                pages_text = extract_text_by_extension(fname, file_bytes)
            
            if not pages_text:
                st.warning(f"⚠️ No text found in {fname} - skipping")
                continue
            
            # Process each page
            for page_num, page_text in pages_text:
                chunks = intelligent_chunk_text(page_text, max_chunk_size=1000, overlap=200)
                
                if not chunks:
                    continue
                
                # Generate embeddings
                with st.spinner(f"🧠 Generating embeddings for page {page_num} of {fname}..."):
                    embeddings = generate_embeddings_batch(openai_client, chunks)
                
                # Create documents for upload
                for cidx, (chunk, emb) in enumerate(zip(chunks, embeddings)):
                    docs_to_upload.append({
                        "id": f"{fid}_p{page_num}_c{cidx}",
                        "content": chunk,
                        "content_vector": emb,
                        "page_number": page_num,
                        "file_name": fname,
                        "chunk_id": f"page_{page_num}_chunk_{cidx}",
                        "document_type": os.path.splitext(fname)[1].lower().lstrip('.'),
                        "last_modified": file_meta.get("lastModifiedDateTime", "2024-01-01T00:00:00Z")
                    })
            
            successful_files = 1
            st.success(f"✅ Successfully processed {fname}")
            
        except Exception as e:
            logger.error("Error processing %s: %s", fname, e)
            st.error(f"❌ Error processing {fname}: {e}")
        
        # Update progress and metrics
        progress = i / total_files
        progress_bar.progress(progress)
        
        processed_files_metric.metric("✅ Files Processed", f"{i}/{total_files}")
        chunks_generated_metric.metric("📄 Chunks Generated", len(docs_to_upload))
        success_rate = (successful_files / i) * 100
        success_rate_metric.metric("📈 Success Rate", f"{success_rate:.1f}%")
    
    # ────────── UPLOAD TO INDEX ──────────
    if not docs_to_upload:
        st.warning("❌ No document chunks generated — nothing to upload.")
        st.stop()
    
    st.markdown("---")
    st.markdown("### 📤 Uploading to Search Index")
    st.info(f"⬆️ Uploading {len(docs_to_upload)} chunks to Azure AI Search...")
    
    with st.spinner("Uploading to search index..."):
        try:
            success = upload_to_search_index_batch(index_name, docs_to_upload, credential)
            
            if success:
                st.success("🎉 All documents uploaded successfully!")
                st.session_state["index_name"] = index_name
                st.session_state["total_documents"] = len(docs_to_upload)
                
                # Final summary with enhanced styling
                st.markdown("---")
                st.markdown('<h3 class="section-header">🎯 Processing Complete!</h3>', unsafe_allow_html=True)
                
                summary_cols = st.columns(4)
                summary_cols[0].metric("📄 Files Processed", successful_files)
                summary_cols[1].metric("🔍 Chunks Created", len(docs_to_upload))
                summary_cols[2].metric("📈 Success Rate", f"{(successful_files/total_files)*100:.1f}%")
                summary_cols[3].metric("🏷️ Index Name", index_name)
                
            
                st.info("🎮 Navigate to the **Chat** tab to start asking questions about your documents!")
                
                # Clear selections after successful processing
                col1, col2 = st.columns(2)
                with col1:
                    if st.button("🔄 Clear Selections & Process More Files", use_container_width=True):
                        st.session_state.file_selections = {}
                        st.session_state.folder_file_mapping = {}
                        st.rerun()
                
                with col2:
                    if st.button("📊 View Processing Summary", use_container_width=True):
                        st.json({
                            "total_files_processed": successful_files,
                            "total_chunks_created": len(docs_to_upload),
                            "success_rate": f"{(successful_files/total_files)*100:.1f}%",
                            "index_name": index_name
                        })
                
            else:
                st.error("❌ Some documents failed to upload to the search index.")
                
        except Exception as e:
            st.error(f"❌ Error uploading to search index: {e}")
            logger.error(f"Upload failed: {e}")


def load_pdf_file():

    uploaded_file = st.file_uploader(
        "Upload file",
        type=["pdf", "docx", "pptx", "csv", "xlsx", "xls"]
    )

    return uploaded_file
def sanitize_document_key(key: str) -> str:
    """
    Sanitize the key to only contain allowed characters.
    A simple approach could be replacing unwanted characters with underscores or encoding using Base64.
    """
    # Replace spaces and unwanted characters with underscores
    sanitized_key = key.replace(' ', '_').replace('.', '_')
    
    # Only keep allowed characters (optional)
    allowed_chars = set("abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789-_=")
    sanitized_key = ''.join(c if c in allowed_chars else '_' for c in sanitized_key)
    
    # If stricter sanitation is needed, you can encode using Base64
    base64_encoded_key = base64.urlsafe_b64encode(sanitized_key.encode()).decode().rstrip('=')
    
    return base64_encoded_key  # Return base64-encoded key to ensure safety
def upload_main():
    
    # Prefer SSO-derived email if available
    default_email = st.session_state.get("user_email", "")
    email_id = st.text_input("Your Email", value=default_email, disabled=bool(default_email))

    # # Optional: support username/password authentication for upload as requested
    # with st.expander("Use username/password for upload (optional)"):
    #     username = st.text_input("Username", key="basic_user")
    #     password = st.text_input("Password", type="password", key="basic_pass")
    #     if username and password:
    #         st.session_state["upload_basic_auth"] = {"username": username, "password": password}
    #         st.info("Basic credentials captured for upload session.")

    pdf_files = load_pdf_file()

    # Only save uploaded file and email to session state, do NOT process here
    if pdf_files is not None and (email_id.strip() or st.session_state.get("user_email")):
        st.session_state["uploaded_pdf_file"] = pdf_files
        st.session_state["uploaded_email_id"] = email_id or st.session_state.get("user_email", "")
        st.success("File uploaded and saved! You can now go to the Data Pipeline tab to run ingestion.")
    elif pdf_files is not None:
        st.warning("Email required: please sign in or enter your email above.")
    else:
        st.info("Please upload a supported file (pdf, docx, pptx, csv, xlsx, xls).")

def data_pipeline_upload():
    import streamlit as st
    # Retrieve uploaded file and email from session state
    pdf_files = st.session_state.get("uploaded_pdf_file")
    email_id = st.session_state.get("uploaded_email_id", "")
    if pdf_files is None or not email_id.strip():
        st.error("No uploaded file or email found. Please upload a file in the Login tab first.")
        return

    # Ingestion/indexing logic moved here
    #credential = DefaultAzureCredential()
    # credential = AzureKeyCredential(OPENAI_API_KEY)
    # openai_client = AzureOpenAI(
    #     azure_ad_token_provider=lambda: credential
    #         .get_token("https://cognitiveservices.azure.com/.default").token,
    #      azure_endpoint=OPENAI_ENDPOINT_URL,
    #      api_version="2024-02-01"
    #  )
    
    openai_client = get_openai_client()
    site_name = email_id.split("@")[0] if "@" in email_id else email_id
    index_name = f"{site_name.lower()}-enhanced-rag-index"

    st.info(f"🏗️ Creating enhanced Azure Search index: **{index_name}**")
    with st.spinner("Creating index with hybrid search and semantic ranking..."):
        create_enhanced_search_index(index_name)
    st.success("✅ Enhanced index created successfully!")

    docs_to_upload = []
    total_files = 1
    progress_bar = st.progress(0)
    status_text = st.empty()
    file_name = pdf_files.name
    file_id = file_name
    status_text.text(f"📄 Processing {file_name}")
    try:
        file_bytes = pdf_files.read()
        pages_text = extract_text_by_extension(file_name, file_bytes)
        for page_num, page_text in pages_text:
            if not page_text.strip():
                continue
            chunks = intelligent_chunk_text(page_text, max_chunk_size=1000, overlap=200)
            if not chunks:
                continue
            embeddings = generate_embeddings_batch(openai_client, chunks)
            for chunk_idx, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                doc_id = sanitize_document_key(f"{file_id}_p{page_num}_c{chunk_idx}")
                doc = {
                    "id": doc_id,
                    "content": chunk,
                    "content_vector": embedding,
                    "page_number": page_num,
                    "file_name": file_name,
                    "chunk_id": f"page_{page_num}_chunk_{chunk_idx}",
                    "document_type": os.path.splitext(file_name)[1].lower().lstrip('.'),
                }
                docs_to_upload.append(doc)
        progress_bar.progress(total_files)
    except Exception as e:
        st.warning(f"⚠️ Error processing {file_name}: {str(e)}")
        logger.error(f"Error processing file {file_name}: {str(e)}")
    progress_bar.progress(1.0)
    status_text.text(f"📊 Generated {len(docs_to_upload)} document chunks")
    if docs_to_upload:
        st.info(f"⬆️ Uploading {len(docs_to_upload)} chunks to Azure AI Search...")
        with st.spinner("Uploading documents..."):
            #credential = DefaultAzureCredential()
            credential = AzureKeyCredential(SEARCH_API_KEY)
            upload_success = upload_to_search_index_batch(index_name, docs_to_upload, credential)
        if upload_success:
            st.success("🎉 Enhanced indexing completed successfully!")
            st.session_state['index_name'] = index_name
            st.session_state['total_documents'] = len(docs_to_upload)
            st.markdown("### 📊 Indexing Summary")
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("📄 Files Processed", "1")
            with col2:
                st.metric("🔍 Chunks Created", len(docs_to_upload))
            with col3:
                st.metric("🎯 Success Rate", "100%")
            st.markdown("### 🚀 Ready for Enhanced Chat!")
            st.info("Navigate to the **Chat** tab to start asking questions with hybrid search and semantic ranking!")
            # Clear uploaded file and email from session state after successful ingestion
            st.session_state.pop("uploaded_pdf_file", None)
            st.session_state.pop("uploaded_email_id", None)
        else:
            st.error("❌ Error uploading documents to search index.")
    else:
        st.warning("⚠️ No documents were processed successfully.")



if __name__ == "__main__":
    main()

